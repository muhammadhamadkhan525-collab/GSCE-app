import streamlit as st
import google.generativeai as genai
from pptx import Presentation
import json
import os
import tempfile
from fpdf import FPDF

# --- CONFIGURATION ---
GEMINI_API_KEY = st.secrets["GEMINI_API_KEY"]
MEMORY_FILE = "quiz_memory.json"

genai.configure(api_key=GEMINI_API_KEY)

# --- MEMORY FUNCTIONS ---
def load_memory():
    """Loads saved quizzes and scores from a local JSON file."""
    if os.path.exists(MEMORY_FILE):
        with open(MEMORY_FILE, "r") as file:
            return json.load(file)
    return {"quizzes": {}}

def save_memory(data):
    """Saves quizzes and scores to a local JSON file."""
    with open(MEMORY_FILE, "w") as file:
        json.dump(data, file, indent=4)

# --- PDF EXPORT FUNCTION ---
def generate_pdf_report(db):
    """Generates a PDF file of all quizzes and returns it as bytes."""
    pdf = FPDF()
    for name, info in db["quizzes"].items():
        pdf.add_page()
        pdf.set_font("Arial", 'B', 16)
        
        # Replace non-standard characters to avoid PDF errors
        safe_name = name.encode('latin-1', 'replace').decode('latin-1')
        pdf.cell(0, 10, f"Quiz: {safe_name}", ln=True)
        
        pdf.set_font("Arial", '', 12)
        pdf.cell(0, 10, f"Attempts: {info['attempts']} | Best Score: {info['best_score']:.1f}%", ln=True)
        pdf.ln(5)
        
        questions = info["data"].get("questions", [])
        for i, q in enumerate(questions):
            # Question Text
            pdf.set_font("Arial", 'B', 12)
            q_text = f"Q{i+1}: {q['question']}".encode('latin-1', 'replace').decode('latin-1')
            pdf.multi_cell(0, 8, q_text)
            
            # Options Text
            pdf.set_font("Arial", '', 12)
            for opt in q['options']:
                opt_text = f" - {opt}".encode('latin-1', 'replace').decode('latin-1')
                pdf.multi_cell(0, 8, opt_text)
                
            # Answer Text
            pdf.set_font("Arial", 'I', 12)
            ans_text = f"Answer: {q['answer']}".encode('latin-1', 'replace').decode('latin-1')
            pdf.multi_cell(0, 8, ans_text)
            pdf.ln(5)
            
    # Save the PDF to a temporary file, read it as bytes, and clean it up
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        pdf.output(tmp.name)
        with open(tmp.name, "rb") as f:
            pdf_bytes = f.read()
    os.remove(tmp.name)
    
    return pdf_bytes

# --- CORE FUNCTIONS ---
def extract_text_from_ppt(ppt_file):
    """Extracts text from an uploaded PPTX file."""
    prs = Presentation(ppt_file)
    text_content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())
    return "\n".join(text_content)

def generate_quiz(text, num_questions):
    """Calls Gemini AI to generate a JSON formatted quiz with a dynamic number of questions."""
    model = genai.GenerativeModel('gemini-2.5-flash')
    prompt = f"""
    Act as a GCSE examiner. Create a {num_questions}-question multiple-choice quiz based on the following text.
    Return ONLY a raw JSON object. Do not include markdown formatting like ```json.
    Structure:
    {{
      "questions": [
        {{
          "question": "The question text?",
          "options": ["Option A", "Option B", "Option C", "Option D"],
          "answer": "Option A",
          "hint": "A helpful hint"
        }}
      ]
    }}
    
    Text: {text}
    """
    response = model.generate_content(prompt)
    
    # Clean up the response in case the AI adds markdown blocks
    raw_text = response.text.strip().replace("```json", "").replace("```", "")
    return json.loads(raw_text)

# --- STREAMLIT USER INTERFACE ---
st.set_page_config(page_title="AI Quiz Master", layout="wide")
st.title("🧠 AI Quiz Master")

# Initialize memory in session state
if "db" not in st.session_state:
    st.session_state.db = load_memory()

# Create tabs for navigation
tab1, tab2, tab3 = st.tabs(["Upload & Generate", "Take a Quiz", "View History & Download"])

# --- TAB 1: GENERATE QUIZ ---
with tab1:
    st.header("Generate a New Quiz")
    quiz_name = st.text_input("Give this quiz a name (e.g., Biology Chapter 1):")
    
    # NEW FEATURE 1: User controls the number of questions
    num_questions = st.number_input("How many questions should the quiz have?", min_value=1, max_value=20, value=5)
    
    uploaded_file = st.file_uploader("Upload your PowerPoint (.pptx)", type="pptx")
    
    if st.button("Generate Quiz") and uploaded_file and quiz_name:
        with st.spinner(f"Extracting text and generating {num_questions} questions..."):
            try:
                slide_text = extract_text_from_ppt(uploaded_file)
                quiz_data = generate_quiz(slide_text, num_questions)
                
                # Save to our "database"
                st.session_state.db["quizzes"][quiz_name] = {
                    "data": quiz_data,
                    "attempts": 0,
                    "best_score": 0
                }
                save_memory(st.session_state.db)
                st.success(f"Quiz '{quiz_name}' with {num_questions} questions generated successfully!")
            except Exception as e:
                st.error(f"Something went wrong: {e}")

# --- TAB 2: TAKE QUIZ ---
with tab2:
    st.header("Take a Saved Quiz")
    saved_quizzes = list(st.session_state.db["quizzes"].keys())
    
    if not saved_quizzes:
        st.info("No quizzes generated yet. Go to the first tab to make one!")
    else:
        selected_quiz = st.selectbox("Select a quiz to take:", saved_quizzes)
        
        if selected_quiz:
            quiz_content = st.session_state.db["quizzes"][selected_quiz]["data"]
            questions = quiz_content.get("questions", [])
            
            with st.form(key=f"quiz_form_{selected_quiz}"):
                user_answers = {}
                for i, q in enumerate(questions):
                    st.write(f"**Q{i+1}: {q['question']}**")
                    with st.expander("Need a hint?"):
                        st.write(q['hint'])
                    
                    user_answers[i] = st.radio(
                        "Select an answer:", 
                        q['options'], 
                        key=f"q_{i}",
                        index=None
                    )
                    st.divider()
                
                submitted = st.form_submit_button("Submit Answers")
                
                if submitted:
                    score = 0
                    for i, q in enumerate(questions):
                        if user_answers[i] == q['answer']:
                            score += 1
                            
                    percentage = (score / len(questions)) * 100
                    st.success(f"You scored {score}/{len(questions)} ({percentage:.1f}%)!")
                    
                    # Update memory with results
                    db_entry = st.session_state.db["quizzes"][selected_quiz]
                    db_entry["attempts"] += 1
                    if percentage > db_entry["best_score"]:
                        db_entry["best_score"] = percentage
                    
                    save_memory(st.session_state.db)

# --- TAB 3: VIEW HISTORY & DOWNLOAD ---
with tab3:
    st.header("Your Learning Dashboard")
    if not st.session_state.db["quizzes"]:
        st.info("No data available yet.")
    else:
        # THE FIX: Two-step download to prevent auto-triggering
        st.subheader("Export")
        if st.button("Prepare PDF for Download"):
            with st.spinner("Building your PDF..."):
                st.session_state.pdf_data = generate_pdf_report(st.session_state.db)
                
        # Only show the actual download button if the PDF has been built
        if "pdf_data" in st.session_state:
            st.download_button(
                label="📄 Click Here to Save PDF",
                data=st.session_state.pdf_data,
                file_name="My_Quiz_History.pdf",
                mime="application/pdf"
            )
            
        st.divider()

        # Show the whole quiz in history
        st.subheader("Quiz History")
        for name, info in st.session_state.db["quizzes"].items():
            with st.expander(f"📖 {name} (Best Score: {info['best_score']:.1f}% | Attempts: {info['attempts']})"):
                questions = info["data"].get("questions", [])
                for i, q in enumerate(questions):
                    st.write(f"**Q{i+1}: {q['question']}**")
                    for opt in q['options']:
                        st.write(f"- {opt}")
                    st.success(f"**Correct Answer:** {q['answer']}")
                    st.divider()